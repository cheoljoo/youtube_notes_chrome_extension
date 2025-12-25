#!/usr/bin/env python3
"""
YouTube Notes Chrome Extension - Development & Publishing Process PPTX Generator
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 31, 31)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

def add_content_slide(title, content_list):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(211, 47, 47)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_before = Pt(8)
        p.space_after = Pt(8)

# Slide 1: Title
add_title_slide("YouTube Notes", "Chrome Extension Development & Publishing")

# Slide 2: Project Overview
add_content_slide("Project Overview", [
    "📌 Extension Name: YouTube Notes",
    "✨ Purpose: Save YouTube video info with tags, opinions & export as CSV",
    "🎯 Target Users: Students, researchers, content creators",
    "📊 Key Features: Auto metadata extraction, tag management, CSV export"
])

# Slide 3: Development Phases
add_content_slide("Development Phases", [
    "Phase 1️⃣: Core Functionality Setup (popup.html, popup.js)",
    "Phase 2️⃣: Storage & Data Management (chrome.storage API)",
    "Phase 3️⃣: YouTube Metadata Extraction (scripting API)",
    "Phase 4️⃣: CSV Export Feature",
    "Phase 5️⃣: UI/UX Enhancement (filtering, tag management)"
])

# Slide 4: Key Features
add_content_slide("Key Features Implemented", [
    "✅ YouTube Title & Publish Date Extraction",
    "✅ Smart Tag Management (reusable tags with suggestions)",
    "✅ Opinion Notes (multi-line textarea)",
    "✅ Tag-based Filtering (click tag to filter notes)",
    "✅ CSV Download (with all metadata)",
    "✅ Local-only Storage (privacy-first)"
])

# Slide 5: Project Structure
add_content_slide("Project File Structure", [
    "📁 manifest.json - Extension configuration & permissions",
    "📁 popup.html/js - Main UI & user interactions",
    "📁 options.html/js - Settings page",
    "📁 background.js - Service worker",
    "📁 images/ - Icons (16x16, 48x48, 128x128 PNG)",
    "📁 privacy_policy.html - Privacy documentation"
])

# Slide 6: Development Tools & Technologies
add_content_slide("Development Tools", [
    "💻 Language: JavaScript (ES6+)",
    "🛠️ APIs: Chrome Storage, Chrome Scripting, Chrome Tabs",
    "🎨 UI: HTML5, CSS3",
    "📦 Build: Python (for icon generation)",
    "📊 Data Format: JSON (storage), CSV (export)"
])

# Slide 7: Testing & Quality Assurance
add_content_slide("Testing & QA Process", [
    "1️⃣ Load extension in development mode (chrome://extensions)",
    "2️⃣ Test on real YouTube videos",
    "3️⃣ Verify metadata extraction accuracy",
    "4️⃣ Check tag filtering functionality",
    "5️⃣ Validate CSV export format",
    "6️⃣ Test on multiple YouTube pages"
])

# Slide 8: Permission Justification
add_content_slide("Permission Justification", [
    "🔒 storage - Local note storage (no cloud)",
    "📑 activeTab - Current tab URL capture",
    "⚙️ scripting - YouTube metadata extraction",
    "🌐 Host Permission - YouTube.com access only",
    "✨ NO identity/OAuth - Local storage only!"
])

# Slide 9: Pre-Submission Checklist
add_content_slide("Pre-Submission Checklist", [
    "✅ manifest.json properly configured",
    "✅ Icons created (16, 48, 128 PNG)",
    "✅ Privacy policy written (English)",
    "✅ README.md with usage instructions",
    "✅ All files packaged as .zip",
    "✅ No external code/resources"
])

# Slide 10: Chrome Web Store Setup
add_content_slide("Chrome Web Store - Account Setup", [
    "1️⃣ Visit: chrome.google.com/webstore/devconsole",
    "2️⃣ Sign in with Google account",
    "3️⃣ Pay $5 USD developer registration fee",
    "4️⃣ Go to Account tab → Add contact email",
    "5️⃣ Verify email address (check inbox)",
    "6️⃣ Ready for submission!"
])

# Slide 11: Submission Process
add_content_slide("Submission Process", [
    "1️⃣ Create new item → Upload youtube_notes_v1.0.0.zip",
    "2️⃣ Fill basic info: Name, Description, Category",
    "3️⃣ Upload icon (128x128) & promo tiles (440x280)",
    "4️⃣ Add 1-2 screenshots (1280x800 recommended)",
    "5️⃣ Set language (English, Korean, etc)",
    "6️⃣ Save Draft → Click Publish"
])

# Slide 12: Detailed Submission Form
add_content_slide("Store Listing Details", [
    "📝 Title: YouTube Notes",
    "📝 Short Description: Save YouTube info with tags",
    "📝 Detailed Description: Full feature overview",
    "📝 Category: Productivity",
    "📝 Language: English (also Korean for international)",
    "📝 Support email: Your contact email"
])

# Slide 13: Privacy & Compliance
add_content_slide("Privacy & Compliance Declaration", [
    "✅ No remote code usage",
    "✅ No user data collection (local only)",
    "✅ No third-party data sharing",
    "✅ Honest permission requests",
    "✅ Clear privacy policy provided",
    "✅ All 3 policy agreements checked"
])

# Slide 14: Review Process
add_content_slide("Chrome Web Store Review", [
    "⏳ Review Duration: 1-3 days (usually)",
    "🔍 Automated Checks: Malware, code analysis",
    "👥 Manual Review: Functionality & policy compliance",
    "⚠️ Host permission = In-depth review (but approved if justified)",
    "📧 Notification via email when approved/rejected"
])

# Slide 15: After Approval
add_content_slide("After Approval ✅", [
    "🎉 Extension goes live on Chrome Web Store",
    "🔍 Users can search: 'YouTube Notes'",
    "⭐ Users can rate & review your extension",
    "📊 View analytics: Users, ratings, crashes",
    "🔄 Update extension: Upload new .zip, increase version"
])

# Slide 16: Updates & Maintenance
add_content_slide("Updates & Maintenance", [
    "1️⃣ Increase version in manifest.json (e.g., 1.0.1)",
    "2️⃣ Update files (popup.js, etc)",
    "3️⃣ Re-zip and upload to store",
    "4️⃣ Write release notes for update",
    "5️⃣ Monitor user reviews & feedback",
    "6️⃣ Plan next features based on feedback"
])

# Slide 17: Best Practices
add_content_slide("Best Practices Going Forward", [
    "🛡️ Security: Never store sensitive data externally",
    "⚡ Performance: Minimize background processes",
    "👤 UX: Keep interface simple & intuitive",
    "📱 Compatibility: Test on different Chrome versions",
    "📢 Marketing: Promote on GitHub, forums, Twitter",
    "💬 Support: Respond to user issues promptly"
])

# Slide 18: Troubleshooting
add_content_slide("Common Issues & Solutions", [
    "❌ 'Unable to publish' → Check Account tab email",
    "❌ 'Invalid manifest' → Validate manifest.json syntax",
    "❌ Icon issues → Ensure PNG 16x16, 48x48, 128x128",
    "❌ Rejection → Review error message, re-submit",
    "❌ Slow review → Check for policy violations"
])

# Slide 19: Success Timeline
add_content_slide("Project Success Timeline", [
    "📅 Days 1-3: Core development & testing",
    "📅 Days 4-5: Polish UI & documentation",
    "📅 Day 6: Package & prepare submission",
    "📅 Day 7: Submit to Chrome Web Store",
    "📅 Days 8-10: Review & approval",
    "📅 Day 11: Live on Chrome Web Store! 🚀"
])

# Slide 20: Final Summary
add_title_slide("Congratulations!", "Your Chrome Extension is Live! 🎉")

# Save presentation
output_path = 'YouTube_Notes_Development_Guide.pptx'
prs.save(output_path)
print(f"✅ PPTX file created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
